"""文档读取模块：docx 分块读取 + xlsx/xls + pdf + pptx + 图片。"""
from __future__ import annotations
import base64
import json
from pathlib import Path
from typing import Any

from config import PROJECTS_DIR, CHUNK_SIZE_PARAGRAPHS
import config


# ---------------------------------------------------------------------------
# 路径工具
# ---------------------------------------------------------------------------

def _resolve(project_name: str, filename: str) -> Path:
    path = PROJECTS_DIR / project_name / filename
    # 安全检查：防止路径穿越
    path = path.resolve()
    base = PROJECTS_DIR.resolve()
    if not str(path).startswith(str(base)):
        raise ValueError(f"路径不合法: {path}")
    return path


# ---------------------------------------------------------------------------
# list_projects / list_files
# ---------------------------------------------------------------------------

def list_projects() -> str:
    if not PROJECTS_DIR.exists():
        return json.dumps({"error": "projects/ 目录不存在"}, ensure_ascii=False)
    projects = [p.name for p in sorted(PROJECTS_DIR.iterdir()) if p.is_dir()]
    return json.dumps({"projects": projects}, ensure_ascii=False)


SUPPORTED_EXTENSIONS = {
    ".docx", ".xlsx", ".xls",
    ".pdf",
    ".pptx", ".ppt",
    ".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif",
}


def list_files(project_name: str) -> str:
    project_dir = (PROJECTS_DIR / project_name).resolve()
    if not project_dir.exists():
        return json.dumps({"error": f"项目 '{project_name}' 不存在"}, ensure_ascii=False)
    files = []
    for f in sorted(project_dir.rglob("*")):
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS:
            size_mb = round(f.stat().st_size / 1024 / 1024, 2)
            rel = f.relative_to(project_dir)
            files.append({"name": str(rel), "size_mb": size_mb, "type": f.suffix.lower().lstrip(".")})
    return json.dumps({"project": project_name, "files": files}, ensure_ascii=False)


# ---------------------------------------------------------------------------
# read_docx
# ---------------------------------------------------------------------------

def read_docx(project_name: str, filename: str, chunk_index: int = 0) -> str:
    try:
        from docx import Document
    except ImportError:
        return json.dumps({"error": "python-docx 未安装"}, ensure_ascii=False)

    path = _resolve(project_name, filename)
    if not path.exists():
        return json.dumps({"error": f"文件不存在: {filename}"}, ensure_ascii=False)

    try:
        doc = Document(str(path))
    except Exception as e:
        return json.dumps({"error": f"无法打开文档: {e}"}, ensure_ascii=False)

    paragraphs = [p.text for p in doc.paragraphs]
    total = len(paragraphs)
    total_chunks = max(1, (total + CHUNK_SIZE_PARAGRAPHS - 1) // CHUNK_SIZE_PARAGRAPHS)

    if chunk_index < 0 or chunk_index >= total_chunks:
        return json.dumps(
            {"error": f"chunk_index 超出范围 [0, {total_chunks - 1}]"},
            ensure_ascii=False,
        )

    start = chunk_index * CHUNK_SIZE_PARAGRAPHS
    end = min(start + CHUNK_SIZE_PARAGRAPHS, total)
    chunk = paragraphs[start:end]

    lines = [f"[{start + i}] {text}" for i, text in enumerate(chunk)]

    return json.dumps(
        {
            "project": project_name,
            "filename": filename,
            "chunk_index": chunk_index,
            "total_chunks": total_chunks,
            "paragraph_range": [start, end - 1],
            "content": "\n".join(lines),
        },
        ensure_ascii=False,
    )


# ---------------------------------------------------------------------------
# read_excel
# ---------------------------------------------------------------------------

def read_excel(project_name: str, filename: str, sheet_name: str = "") -> str:
    path = _resolve(project_name, filename)
    if not path.exists():
        return json.dumps({"error": f"文件不存在: {filename}"}, ensure_ascii=False)

    suffix = path.suffix.lower()

    if suffix == ".xlsx":
        return _read_xlsx(path, sheet_name)
    elif suffix == ".xls":
        return _read_xls(path, sheet_name)
    else:
        return json.dumps({"error": f"不支持的格式: {suffix}"}, ensure_ascii=False)


def _read_xlsx(path: Path, sheet_name: str) -> str:
    try:
        import openpyxl
    except ImportError:
        return json.dumps({"error": "openpyxl 未安装"}, ensure_ascii=False)

    try:
        wb = openpyxl.load_workbook(str(path), data_only=True)
    except Exception as e:
        return json.dumps({"error": f"无法打开文件: {e}"}, ensure_ascii=False)

    sheets_to_read = [sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.sheetnames
    result: dict[str, Any] = {}

    for sname in sheets_to_read:
        ws = wb[sname]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(c) if c is not None else "" for c in row])
        result[sname] = rows

    return json.dumps(
        {"filename": path.name, "sheets": result},
        ensure_ascii=False,
    )


def _read_xls(path: Path, sheet_name: str) -> str:
    try:
        import xlrd
    except ImportError:
        return json.dumps({"error": "xlrd 未安装"}, ensure_ascii=False)

    try:
        wb = xlrd.open_workbook(str(path))
    except Exception as e:
        return json.dumps({"error": f"无法打开文件: {e}"}, ensure_ascii=False)

    sheet_names = wb.sheet_names()
    sheets_to_read = [sheet_name] if sheet_name and sheet_name in sheet_names else sheet_names
    result: dict[str, Any] = {}

    for sname in sheets_to_read:
        ws = wb.sheet_by_name(sname)
        rows = []
        for r in range(ws.nrows):
            rows.append([str(ws.cell_value(r, c)) for c in range(ws.ncols)])
        result[sname] = rows

    return json.dumps(
        {"filename": path.name, "sheets": result, "note": ".xls 为只读格式"},
        ensure_ascii=False,
    )


# ---------------------------------------------------------------------------
# search_in_file
# ---------------------------------------------------------------------------

def search_in_file(project_name: str, filename: str, keyword: str) -> str:
    path = _resolve(project_name, filename)
    if not path.exists():
        return json.dumps({"error": f"文件不存在: {filename}"}, ensure_ascii=False)

    suffix = path.suffix.lower()

    if suffix == ".docx":
        return _search_docx(path, keyword)
    elif suffix in (".xlsx", ".xls"):
        return _search_excel(path, keyword, suffix)
    else:
        return json.dumps({"error": f"不支持的格式: {suffix}"}, ensure_ascii=False)


def _search_docx(path: Path, keyword: str) -> str:
    try:
        from docx import Document
    except ImportError:
        return json.dumps({"error": "python-docx 未安装"}, ensure_ascii=False)

    doc = Document(str(path))
    matches = []
    for i, para in enumerate(doc.paragraphs):
        if keyword in para.text:
            matches.append({"paragraph_index": i, "text": para.text})

    return json.dumps(
        {"filename": path.name, "keyword": keyword, "matches": matches},
        ensure_ascii=False,
    )


def _search_excel(path: Path, keyword: str, suffix: str) -> str:
    matches = []

    if suffix == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)
            for sname in wb.sheetnames:
                ws = wb[sname]
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None and keyword in str(cell.value):
                            matches.append({
                                "sheet": sname,
                                "cell": cell.coordinate,
                                "value": str(cell.value),
                            })
        except Exception as e:
            return json.dumps({"error": str(e)}, ensure_ascii=False)
    else:
        try:
            import xlrd
            wb = xlrd.open_workbook(str(path))
            for sname in wb.sheet_names():
                ws = wb.sheet_by_name(sname)
                for r in range(ws.nrows):
                    for c in range(ws.ncols):
                        val = str(ws.cell_value(r, c))
                        if keyword in val:
                            col_letter = chr(ord("A") + c) if c < 26 else f"C{c}"
                            matches.append({
                                "sheet": sname,
                                "cell": f"{col_letter}{r + 1}",
                                "value": val,
                            })
        except Exception as e:
            return json.dumps({"error": str(e)}, ensure_ascii=False)

    return json.dumps(
        {"filename": path.name, "keyword": keyword, "matches": matches},
        ensure_ascii=False,
    )


# ---------------------------------------------------------------------------
# read_image
# ---------------------------------------------------------------------------

_MIME_MAP = {
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png",
    ".bmp": "image/bmp", ".webp": "image/webp", ".gif": "image/gif",
}


def read_image(project_name: str, filename: str, question: str = "") -> str:
    path = _resolve(project_name, filename)
    if not path.exists():
        return json.dumps({"error": f"文件不存在: {filename}"}, ensure_ascii=False)

    mime = _MIME_MAP.get(path.suffix.lower())
    if not mime:
        return json.dumps({"error": f"不支持的图片格式: {path.suffix}"}, ensure_ascii=False)

    try:
        b64 = base64.b64encode(path.read_bytes()).decode()
    except Exception as e:
        return json.dumps({"error": f"读取图片失败: {e}"}, ensure_ascii=False)

    data_url = f"data:{mime};base64,{b64}"
    prompt = question or "请详细描述这张图片的内容"

    try:
        from openai import OpenAI
        client = OpenAI(api_key=config.API_KEY, base_url=config.BASE_URL)
        resp = client.chat.completions.create(
            model=config.IMAGE_MODEL,
            messages=[{"role": "user", "content": [
                {"type": "image_url", "image_url": {"url": data_url}},
                {"type": "text", "text": prompt},
            ]}],
            temperature=config.TEMPERATURE,
        )
        description = resp.choices[0].message.content or ""
    except Exception as e:
        return json.dumps({"error": f"视觉API调用失败: {e}"}, ensure_ascii=False)

    return json.dumps(
        {"filename": filename, "question": prompt, "description": description},
        ensure_ascii=False,
    )


# ---------------------------------------------------------------------------
# read_pdf
# ---------------------------------------------------------------------------

def _parse_range(range_str: str, total: int) -> list[int]:
    """解析页码/幻灯片范围字符串，返回0-based索引列表。"""
    if not range_str.strip():
        return list(range(total))
    indices = []
    for part in range_str.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            start = max(0, int(a.strip()) - 1)
            end = min(total - 1, int(b.strip()) - 1)
            indices.extend(range(start, end + 1))
        elif part.isdigit():
            idx = int(part) - 1
            if 0 <= idx < total:
                indices.append(idx)
    return indices


def read_pdf(project_name: str, filename: str, page_range: str = "") -> str:
    try:
        import fitz  # pymupdf
    except ImportError:
        return json.dumps({"error": "pymupdf 未安装，请运行 pip install pymupdf"}, ensure_ascii=False)

    path = _resolve(project_name, filename)
    if not path.exists():
        return json.dumps({"error": f"文件不存在: {filename}"}, ensure_ascii=False)

    try:
        doc = fitz.open(str(path))
    except Exception as e:
        return json.dumps({"error": f"无法打开PDF: {e}"}, ensure_ascii=False)

    total = doc.page_count
    indices = _parse_range(page_range, total)

    pages = []
    for i in indices:
        text = doc[i].get_text()
        pages.append({"page": i + 1, "text": text.strip()})

    return json.dumps(
        {"filename": filename, "total_pages": total, "pages_returned": len(pages), "pages": pages},
        ensure_ascii=False,
    )


# ---------------------------------------------------------------------------
# read_pptx
# ---------------------------------------------------------------------------

def read_pptx(project_name: str, filename: str, slide_range: str = "") -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return json.dumps({"error": "python-pptx 未安装，请运行 pip install python-pptx"}, ensure_ascii=False)

    path = _resolve(project_name, filename)
    if not path.exists():
        return json.dumps({"error": f"文件不存在: {filename}"}, ensure_ascii=False)

    try:
        prs = Presentation(str(path))
    except Exception as e:
        return json.dumps({"error": f"无法打开PPTX: {e}"}, ensure_ascii=False)

    total = len(prs.slides)
    indices = _parse_range(slide_range, total)

    slides = []
    for i in indices:
        slide = prs.slides[i]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
        slides.append({"slide": i + 1, "text": "\n".join(texts)})

    return json.dumps(
        {"filename": filename, "total_slides": total, "slides_returned": len(slides), "slides": slides},
        ensure_ascii=False,
    )


# ---------------------------------------------------------------------------
# read_project_folder
# ---------------------------------------------------------------------------

def read_project_folder(
    project_name: str,
    subfolder: str = "",
    include_images: bool = True,
    max_files: int = 20,
) -> str:
    base = (PROJECTS_DIR / project_name).resolve()
    if subfolder:
        base = (base / subfolder).resolve()
    if not base.exists():
        return json.dumps({"error": f"目录不存在: {base}"}, ensure_ascii=False)

    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif"}
    TEXT_EXTS = {".docx", ".xlsx", ".xls", ".pdf", ".pptx", ".ppt"}

    all_files = sorted(f for f in base.rglob("*") if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS)
    if not include_images:
        all_files = [f for f in all_files if f.suffix.lower() not in IMAGE_EXTS]

    all_files = all_files[:max_files]

    summaries = []
    for f in all_files:
        rel = str(f.relative_to(PROJECTS_DIR / project_name))
        suffix = f.suffix.lower()
        try:
            if suffix == ".docx":
                raw = json.loads(read_docx(project_name, rel, 0))
                content = raw.get("content", "")[:2000]
                ftype = "docx"
            elif suffix in (".xlsx", ".xls"):
                raw = json.loads(read_excel(project_name, rel))
                content = json.dumps(raw.get("sheets", {}), ensure_ascii=False)[:2000]
                ftype = "excel"
            elif suffix == ".pdf":
                raw = json.loads(read_pdf(project_name, rel, ""))
                pages_text = " ".join(p["text"] for p in raw.get("pages", []))
                content = pages_text[:2000]
                ftype = "pdf"
            elif suffix in (".pptx", ".ppt"):
                raw = json.loads(read_pptx(project_name, rel, ""))
                slides_text = " ".join(s["text"] for s in raw.get("slides", []))
                content = slides_text[:2000]
                ftype = "pptx"
            elif suffix in IMAGE_EXTS:
                raw = json.loads(read_image(project_name, rel, "请简要描述这张图片的内容和用途"))
                content = raw.get("description", raw.get("error", ""))
                ftype = "image"
            else:
                continue
        except Exception as e:
            content = f"读取失败: {e}"
            ftype = suffix.lstrip(".")

        summaries.append({"file": rel, "type": ftype, "content": content})

    return json.dumps(
        {"project": project_name, "files_read": len(summaries), "summaries": summaries},
        ensure_ascii=False,
    )
